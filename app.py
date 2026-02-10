from flask import Flask, render_template, request, send_file
from werkzeug.utils import secure_filename
import os, tempfile
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

app = Flask(__name__)

# CONFIG
app.config['UPLOAD_FOLDER'] = tempfile.mkdtemp()
app.config['TEMPLATE_FILE'] = 'shubh.pptx'

ALLOWED_EXTENSIONS = {'png', 'jpg', 'jpeg'}


# ---------- HELPERS ----------

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def clean_text(text):
    return text.replace("\r", "").replace("_x000D_", "").strip()


def set_font_8(paragraph):
    for run in paragraph.runs:
        run.font.size = Pt(8)
        run.font.name = "Segoe UI"


def replace_text_in_shape(shape, placeholder, value):
    if not shape.has_text_frame:
        return

    for paragraph in shape.text_frame.paragraphs:
        if placeholder in paragraph.text:
            paragraph.text = paragraph.text.replace(placeholder, clean_text(value))
            set_font_8(paragraph)


# ---------- ROUTES ----------

@app.route('/')
def index():
    return render_template('index.html')


@app.route('/generate', methods=['POST'])
def generate_ppt():
    try:
        # ---- FORM DATA ----
        project_title = clean_text(request.form['projectTitle'])
        project_domain = clean_text(request.form['projectDomain'])
        guide_name = clean_text(request.form['guideName'])
        brief_idea = clean_text(request.form['briefIdea'])

        # LIMIT BRIEF IDEA TO 100 WORDS
        words = brief_idea.split()
        if len(words) > 100:
            brief_idea = " ".join(words[:100])

        students = [
            request.form.get('student1', ''),
            request.form.get('student2', ''),
            request.form.get('student3', ''),
            request.form.get('student4', '')
        ]

        applications = [
            request.form.get('application1', ''),
            request.form.get('application2', ''),
            request.form.get('application3', ''),
            request.form.get('application4', '')
        ]

        green_image = request.files.get('greenAreaPhoto')
        screenshots = request.files.getlist('projectScreenshots')

        if not green_image or not allowed_file(green_image.filename):
            return "Green area image required.", 400

        if not (1 <= len(screenshots) <= 3):
            return "Upload 1 to 3 screenshots.", 400

        # ---- LOAD PPT ----
        prs = Presentation(app.config['TEMPLATE_FILE'])
        slide = prs.slides[0]

        # ---- TEXT REPLACEMENT ----
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            replace_text_in_shape(shape, "{{PROJECT_TITLE}}", project_title)
            replace_text_in_shape(shape, "{{PROJECT_DOMAIN}}", project_domain)
            replace_text_in_shape(shape, "{{GUIDE_NAME}}", guide_name)
            replace_text_in_shape(shape, "{{PROJECT_DESCRIPTION}}", brief_idea)

            for i in range(4):
                replace_text_in_shape(shape, f"{{{{STUDENT_{i+1}}}}}", students[i])
                replace_text_in_shape(shape, f"{{{{APPLICATION_{i+1}}}}}", applications[i])

        # ---- FIND IMAGE PLACEHOLDERS ----
        green_box = None
        screenshot_boxes = []

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                w, h = shape.width.inches, shape.height.inches

                if 2.3 < w < 2.6 and 1.0 < h < 1.3:
                    green_box = shape

                if 1.4 < w < 1.8 and 0.8 < h < 1.1:
                    screenshot_boxes.append(shape)

        # ---- INSERT GREEN IMAGE ----
        green_path = os.path.join(
            app.config['UPLOAD_FOLDER'],
            secure_filename(green_image.filename)
        )
        green_image.save(green_path)

        slide.shapes.add_picture(
            green_path,
            green_box.left,
            green_box.top,
            green_box.width,
            green_box.height
        )
        slide.shapes._spTree.remove(green_box._element)

        # ---- INSERT SCREENSHOTS ----
        for i, file in enumerate(screenshots):
            if i >= len(screenshot_boxes):
                break

            img_path = os.path.join(
                app.config['UPLOAD_FOLDER'],
                secure_filename(file.filename)
            )
            file.save(img_path)

            box = screenshot_boxes[i]
            slide.shapes.add_picture(
                img_path,
                box.left,
                box.top,
                box.width,
                box.height
            )
            slide.shapes._spTree.remove(box._element)

        output = os.path.join(app.config['UPLOAD_FOLDER'], 'project_presentation.pptx')
        prs.save(output)

        return send_file(output, as_attachment=True)

    except Exception as e:
        return f"Error: {e}", 500


if __name__ == '__main__':
    app.run(debug=True)
